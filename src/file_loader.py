from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation

def load_pdf(file):
    reader = PdfReader(file)
    text = ""

    for page in reader.pages:
        text += page.extract_text() or ""

    return text

def load_docx(file):
    doc = Document(file)
    text = "\n".join([para.text for para in doc.paragraphs])
    return text

def load_pptx(file):
    prs = Presentation(file)
    text = ""

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"

    return text

def load_file(uploaded_file):
    file_name = uploaded_file.name.lower()

    if file_name.endswith(".pdf"):
        return load_pdf(uploaded_file)

    elif file_name.endswith(".docx"):
        return load_docx(uploaded_file)

    elif file_name.endswith(".pptx"):
        return load_pptx(uploaded_file)

    else:
        return ""